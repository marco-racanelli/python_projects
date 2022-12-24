import collections 
import collections.abc
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World 🌎🌍🌏"
subtitle.text = "made with python-pptx"

prs.save('hello_world.pptx')
