import collections 
import collections.abc
from pptx import Presentation
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from datetime import date
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR #,MSO_AUTO_SIZE
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
from pptx.util import Inches

# Create some variables that we'll use later

# ------ Datetime variables
today = date.today()

# ------ Presentation variables
author_1 = 'Marco Racanelli'

# ------ Colors for python-pptx
white = RGBColor(255, 255, 255)
ultramatine_blue = RGBColor(33, 98, 250)
oxford_blue = RGBColor(16, 37, 66)
mikado_yellow = RGBColor(255, 195, 39)
salsa_red = RGBColor(255, 63, 75)

# ------ Colors for plotly
oxford_blue2 = '#102542'
ultramatine_blue2 = '#2162FA'
salsa_red2 = '#FF3F4B'
mikado_yellow2 = '#FFC327'

# Manipulate data and build graphs
df22 = pd.read_excel('revenue22.xlsx')
df21 = pd.read_excel('revenue21.xlsx')
df_pivot22 = pd.pivot_table(df22, index='Quarter', values=['Revenue','Expenses'], aggfunc=np.sum).reset_index()
df_pivot21 = pd.pivot_table(df21, index='Quarter', values=['Revenue','Expenses'], aggfunc=np.sum).reset_index()
fig22 = go.Figure()

fig22.add_trace(go.Bar(x=df_pivot22.Quarter, y=df_pivot22.Expenses,
                base=df_pivot22.Expenses*-1,
                marker_color=salsa_red2,
                name='Expenses',
                #text=round(df_pivot22.Expenses*-1/1000000,2)
                ))
fig22.add_trace(go.Bar(x=df_pivot22.Quarter, y=df_pivot22.Revenue,
                base=0,
                marker_color=ultramatine_blue2,
                name='Revenue',
                #text=round(df_pivot22.Revenue*-1/1000000,2)
                ))

fig22.update_layout(
    title={
        'text':'Revenue VS Expenses',
        'y':0.8,
        'x':0.5,
        'xanchor':'center',
        'yanchor':'top'},
    font=dict(
        family="Arial",
        size=10,
        color=oxford_blue2),
    template='plotly_white',
    autosize=True,
    width=680,
    height=350,
    paper_bgcolor="rgba(0,0,0,0)",
    #plot_bgcolor="rgba(0,0,0,0)"
)
    
fig22.write_image('fig1.png')

profit22 = sum(df22.Revenue)-sum(df22.Expenses)
profit21 = sum(df21.Revenue)-sum(df21.Expenses)
YoY_profit_diff = profit22 - profit21
YoY_profit_diff_perc = str(int(round(YoY_profit_diff / profit21 *100,0)))
YoY_profit_diff_perc

rev21 = pd.pivot_table(df21, index='Category', values='Revenue', aggfunc=np.sum).reset_index()
rev22 = pd.pivot_table(df22, index='Category', values='Revenue', aggfunc=np.sum).reset_index()

categories = rev21.Category

fig2 = go.Figure(data=[
    go.Bar(name='2021', x=categories, y=rev21.Revenue, marker_color= oxford_blue2),
    go.Bar(name='2022', x=categories, y=rev22.Revenue, marker_color= ultramatine_blue2 )
])
# Change the bar mode
fig2.update_layout(
    barmode='group',
    title={
        'text':'YoY comparison',
        'y':0.8,
        'x':0.5,
        'xanchor':'center',
        'yanchor':'top'},
    font=dict(
        family="Arial",
        size=10,
        color=oxford_blue2),
    template='plotly_white',
    autosize=False,
    width=680,
    height=350,
    paper_bgcolor="rgba(0,0,0,0)",
    #plot_bgcolor="rgba(0,0,0,0)"
)
fig2.write_image("fig2.png")

df_profits = pd.read_excel('profits22.xlsx')

profit_pivot = pd.pivot_table(df_profits, index=['Region','Subsidiary'], values='Profit',aggfunc=np.sum).reset_index()

fig3 = px.bar(profit_pivot, x="Region", y="Profit", facet_col="Subsidiary",facet_col_wrap=4,
            category_orders={"Region":profit_pivot.Region.unique(),
                            "Subsidiary": profit_pivot.Subsidiary.unique()},
            facet_row_spacing=0.19, # default is 0.07 when facet_col_wrap is used
            facet_col_spacing=0.02,
            labels={"Region": "","Profit":""}) # default is 0.03)

fig3.for_each_annotation(lambda a: a.update(text=a.text.split("=")[-1]))

fig3.update_layout(
    barmode='group',
    font=dict(
        family="Arial",
        size=10,
        color='black'),
    template='simple_white',
    autosize=False,
    width=680,
    height=350,
    paper_bgcolor="rgba(0,0,0,0)",
    margin=dict(l=50 , r=0, t=90, b=0)
    #plot_bgcolor="rgba(0,0,0,0)"
)
fig3.update_xaxes(showticklabels=True, tickangle=0,tickfont=dict(size=7.5))
fig3.update_traces(width=0.5, marker_color=ultramatine_blue2)
fig3.write_image("fig3.png")

# Observations section: here we store all the observations
observation1 = "We observed a decline in revenue in Q2 and Q3 last year. \nIn Q4, after boosting our advertisement spend, our revenue recovered drastically. \n2022 was overall positive. Our profit grew by " +YoY_profit_diff_perc+"% YoY."
observation2 = "Revenue increased across all our core activities. \nRevenue generated through retailers registered the biggest increase YoY."
observation3 = "Our subsidiaries were profitable across all regions. \nVisio-LAB ltd. had the strongest performance, follwed closely by Farmed-X Inc."

prs = Presentation()

# Title slide ---------------------
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = ultramatine_blue

tf = title.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

run = p.add_run()
run.text = 'Python + pptx'
font = run.font
font.name = 'Arial'
font.size = Pt(60)
font.bold = True
font.color.rgb = white

tfs = subtitle.text_frame
tfs.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tfs.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = 'Presentation model: 2022 yearly recap' 
font = run.font
font.name = 'Arial'
font.size = Pt(24)
font.bold = True
font.color.rgb = white

p00 = tfs.add_paragraph()

p0 = tfs.add_paragraph()
p0.alignment = PP_ALIGN.CENTER
run = p0.add_run()
run.text = 'Author: '+author_1
font = run.font
font.name = 'Arial'
font.size = Pt(16)
font.bold = False
font.color.rgb = white

# Slide 1 ---------------------

# Slide structure
content_slide1 = prs.slide_layouts[1]
slide1 = prs.slides.add_slide(content_slide1)
title1 = slide1.shapes.title
text1 = slide1.placeholders[1]
graph1 = slide1.shapes.add_picture('fig1.png',0,0)
graph1.left = Inches(0)
graph1.top = Inches(2.8)

# Slide look
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = white

tf1 = title1.text_frame
tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
run1 = p1.add_run()
run1.text = '2022 Recap: top-line'
font = run1.font
font.name = 'Arial'
font.size = Pt(36)
font.bold = True
font.color.rgb = oxford_blue

tfs1 = text1.text_frame
tfs1.vertical_anchor = MSO_ANCHOR.TOP
p1 = tfs1.paragraphs[0]
p1.alignment = PP_ALIGN.LEFT
run1 = p1.add_run()
run1.text = observation1
font = run1.font
font.name = 'Arial'
font.size = Pt(16)
font.bold = False
font.color.rgb = oxford_blue

# Slide 2 ---------------------

# Slide structure
content_slide2 = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(content_slide2)
title2 = slide2.shapes.title
text2 = slide2.placeholders[1]
graph2 = slide2.shapes.add_picture('fig2.png',0,0)
graph2.left = Inches(0)
graph2.top = Inches(2.6)

# Slide look
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = white

tf2 = title2.text_frame
tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = '2022 Recap: revenue sources'
font = run2.font
font.name = 'Arial'
font.size = Pt(36)
font.bold = True
font.color.rgb = oxford_blue

tfs2 = text2.text_frame
tfs2.vertical_anchor = MSO_ANCHOR.TOP
p2 = tfs2.paragraphs[0]
p2.alignment = PP_ALIGN.LEFT
run2 = p2.add_run()
run2.text = observation2
font = run2.font
font.name = 'Arial'
font.size = Pt(16)
font.bold = False
font.color.rgb = oxford_blue

# Slide 3 ---------------------

# Slide structure
content_slide3 = prs.slide_layouts[1]
slide3 = prs.slides.add_slide(content_slide2)
title3 = slide3.shapes.title
text3 = slide3.placeholders[1]
graph3 = slide3.shapes.add_picture('fig3.png',0,0)
graph3.left = Inches(0)
graph3.top = Inches(2.6)

# Slide look
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = white

tf3 = title3.text_frame
tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = '2022 Recap: subsidiaries'
font = run3.font
font.name = 'Arial'
font.size = Pt(36)
font.bold = True
font.color.rgb = oxford_blue

tfs3 = text3.text_frame
tfs3.vertical_anchor = MSO_ANCHOR.TOP
p3 = tfs3.paragraphs[0]
p3.alignment = PP_ALIGN.LEFT
run3 = p3.add_run()
run3.text = observation2
font = run3.font
font.name = 'Arial'
font.size = Pt(16)
font.bold = False
font.color.rgb = oxford_blue

# End slide ---------------------
end_slide_layout = prs.slide_layouts[0]
slide_end = prs.slides.add_slide(end_slide_layout)
title_end = slide_end.shapes.title
subtitle_end = slide_end.placeholders[1]

background = slide_end.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = oxford_blue

tf_end = title_end.text_frame
p_end = tf_end.paragraphs[0]
p_end.alignment = PP_ALIGN.CENTER

run_end = p_end.add_run()
run_end.text = 'The end'
font = run_end.font
font.name = 'Arial'
font.size = Pt(50)
font.bold = True
font.color.rgb = white

tfs_end = subtitle_end.text_frame
tfs_end.vertical_anchor = MSO_ANCHOR.MIDDLE
p_end = tfs_end.paragraphs[0]
p_end.alignment = PP_ALIGN.CENTER
run_end = p_end.add_run()
run_end.text = 'Made with 🤍 by Marco Racanelli on '+ today.strftime('%m/%d/%Y')
font = run_end.font
font.name = 'Arial'
font.size = Pt(24)
font.bold = False
font.color.rgb = white

prs.save('presentation.pptx')
